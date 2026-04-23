from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("/Users/wale/Desktop/aiapis/apiweb/AIPIS_PITCH_DECK.pptx")
LOGO = "/Users/wale/Desktop/aiapis/apiweb/logo.png"


def rgb(hex_str):
    hex_str = hex_str.replace("#", "")
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


BLACK = rgb("#0D1120")
DEEP = rgb("#171C30")
WHITE = rgb("#FFFFFF")
MUTED = rgb("#C9D1E5")
SOFT = rgb("#5D657B")
INDIGO = rgb("#625BFF")
CYAN = rgb("#1FB6FF")
CORAL = rgb("#FF7A59")
LIME = rgb("#D6FF72")
ROSE = rgb("#FF5AA5")
PALE = rgb("#F4F7FF")
BORDER = rgb("#DEE4F3")


def add_bg(slide, dark=True):
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLACK if dark else PALE
    bg.line.fill.background()


def add_orb(slide, left, top, size, color, transparency=35):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency / 100
    shape.line.fill.background()


def add_text(slide, left, top, width, height, text, size=20, color=WHITE, bold=False, font="Aptos", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = color
    return box


def add_tag(slide, left, top, text, dark=True):
    w = max(1.35, 0.12 * len(text) + 0.5)
    tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, Inches(w), Inches(0.42))
    tag.fill.solid()
    tag.fill.fore_color.rgb = DEEP if dark else rgb("#EAF0FF")
    tag.line.fill.background()
    add_text(slide, left + Inches(0.18), top + Inches(0.08), Inches(w - 0.2), Inches(0.2), text.upper(), size=10, color=WHITE if dark else BLACK, bold=True)


def add_card(slide, left, top, width, height, title, body, accent, dark=False):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = DEEP if dark else WHITE
    card.line.color.rgb = accent if dark else BORDER
    card.line.width = Pt(1.1)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, Inches(0.08))
    band.fill.solid()
    band.fill.fore_color.rgb = accent
    band.line.fill.background()
    add_text(slide, left + Inches(0.22), top + Inches(0.22), width - Inches(0.35), Inches(0.35), title, size=18, color=WHITE if dark else BLACK, bold=True)
    add_text(slide, left + Inches(0.22), top + Inches(0.68), width - Inches(0.35), height - Inches(0.82), body, size=12, color=MUTED if dark else SOFT)


def cover(slide):
    add_bg(slide, dark=True)
    add_orb(slide, Inches(9.55), Inches(4.65), Inches(2.4), INDIGO, 0)
    add_orb(slide, Inches(10.25), Inches(4.85), Inches(1.9), CYAN, 0)
    add_orb(slide, Inches(10.95), Inches(5.0), Inches(1.35), CORAL, 0)
    slide.shapes.add_picture(LOGO, Inches(0.72), Inches(0.62), height=Inches(0.76))
    add_tag(slide, Inches(0.72), Inches(6.35), "multi-model ai api platform", dark=True)
    add_text(slide, Inches(0.72), Inches(1.45), Inches(5.4), Inches(0.95), "AIPIs", size=34, color=WHITE, bold=True)
    add_text(slide, Inches(0.72), Inches(2.35), Inches(5.9), Inches(2.0), "A unified platform for accessing, routing, and monetizing AI providers like ChatGPT, Claude, Qwen, Ollama, and future model endpoints from one product layer.", size=20, color=MUTED)
    add_text(slide, Inches(0.72), Inches(6.95), Inches(1.0), Inches(0.3), "Pitch Deck", size=13, color=LIME, bold=True)
    add_text(slide, Inches(1.72), Inches(6.95), Inches(1.0), Inches(0.3), "April 2026", size=13, color=MUTED)


def problem(slide):
    add_bg(slide, dark=False)
    add_tag(slide, Inches(0.72), Inches(0.5), "Problem", dark=False)
    add_text(slide, Inches(0.72), Inches(1.0), Inches(8.3), Inches(0.95), "The AI ecosystem is fragmented, fast-moving, and painful to integrate repeatedly.", size=24, color=BLACK, bold=True)
    add_text(slide, Inches(0.72), Inches(1.75), Inches(8.8), Inches(1.1), "Teams want access to the best model for each task, but every provider adds auth, schema differences, pricing logic, fallback behavior, and operational complexity.", size=15, color=SOFT)
    add_card(slide, Inches(0.72), Inches(3.1), Inches(2.75), Inches(1.95), "Too many providers", "OpenAI, Anthropic, Qwen, Ollama, and others all require separate integration and maintenance paths.", CORAL)
    add_card(slide, Inches(3.62), Inches(3.1), Inches(2.75), Inches(1.95), "Switching is expensive", "Moving products between models often means rewriting prompt logic, schemas, and app controls.", CYAN)
    add_card(slide, Inches(6.52), Inches(3.1), Inches(2.75), Inches(1.95), "No shared control", "Without a platform layer, usage, routing, billing, and failover stay fragmented.", INDIGO)
    add_card(slide, Inches(9.42), Inches(3.1), Inches(2.9), Inches(1.95), "Lock-in risk", "One-provider dependency creates business risk the moment model quality, pricing, or availability changes.", LIME)


def product(slide):
    add_bg(slide, dark=True)
    add_tag(slide, Inches(0.72), Inches(0.5), "Product", dark=True)
    add_text(slide, Inches(0.72), Inches(1.0), Inches(8.4), Inches(0.95), "AIPIs gives users and developers one access layer across many AI providers.", size=24, color=WHITE, bold=True)
    add_text(slide, Inches(0.72), Inches(1.75), Inches(7.8), Inches(1.1), "The platform provides a chat experience, provider switchboard, unified API gateway, and control plane for usage, routing, and model orchestration.", size=15, color=MUTED)
    add_card(slide, Inches(0.72), Inches(3.2), Inches(2.75), Inches(1.9), "Chat interface", "Users can interact with different models from one product experience.", CYAN, dark=True)
    add_card(slide, Inches(3.62), Inches(3.2), Inches(2.75), Inches(1.9), "Unified API", "Developers hit one endpoint layer instead of wiring every provider directly.", CORAL, dark=True)
    add_card(slide, Inches(6.52), Inches(3.2), Inches(2.75), Inches(1.9), "Routing engine", "Requests can be sent to the best model by policy, latency, cost, or workload type.", INDIGO, dark=True)
    add_card(slide, Inches(9.42), Inches(3.2), Inches(2.9), Inches(1.9), "Usage controls", "Metering, quotas, and customer plans attach directly to AI usage.", LIME, dark=True)


def architecture(slide):
    add_bg(slide, dark=False)
    add_tag(slide, Inches(0.72), Inches(0.5), "Architecture", dark=False)
    add_text(slide, Inches(0.72), Inches(1.0), Inches(7.8), Inches(0.95), "An orchestration layer between apps and model providers.", size=24, color=BLACK, bold=True)
    add_text(slide, Inches(0.72), Inches(6.2), Inches(10.0), Inches(0.7), "AIPIs can support hosted LLMs, open-weight deployments, local inference, and future providers through one commercial and operational plane.", size=15, color=SOFT)
    steps = [
        ("Input", "User app, chat client, or API consumer", CORAL),
        ("Gateway", "Auth, normalization, quotas", CYAN),
        ("Routing", "Policy, fallback, model selection", INDIGO),
        ("Providers", "Claude, ChatGPT, Qwen, Ollama", ROSE),
        ("Control plane", "Logs, billing, analytics", LIME),
    ]
    x = 0.92
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.3), Inches(3.2), Inches(10.5), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BORDER
    bar.line.fill.background()
    for title, body, accent in steps:
      add_orb(slide, Inches(x + 0.45), Inches(2.95), Inches(0.42), accent, 0)
      add_card(slide, Inches(x), Inches(3.55), Inches(2.0), Inches(1.55), title, body, accent)
      x += 2.28


def moat(slide):
    add_bg(slide, dark=True)
    add_tag(slide, Inches(0.72), Inches(0.5), "Differentiation", dark=True)
    add_text(slide, Inches(0.72), Inches(1.0), Inches(8.8), Inches(0.95), "The moat is control across providers, not dependence on one model vendor.", size=24, color=WHITE, bold=True)
    add_card(slide, Inches(0.72), Inches(3.15), Inches(3.6), Inches(1.95), "Unified switching", "AIPIs makes provider changes operationally cheap for the customer.", CORAL, dark=True)
    add_card(slide, Inches(4.58), Inches(3.15), Inches(3.6), Inches(1.95), "Commercial infrastructure", "Usage, quotas, plans, and product controls are built into the AI layer itself.", CYAN, dark=True)
    add_card(slide, Inches(8.44), Inches(3.15), Inches(3.85), Inches(1.95), "Resilience and speed", "Fallback and routing reduce outages, latency problems, and provider volatility.", INDIGO, dark=True)
    add_text(slide, Inches(0.72), Inches(5.75), Inches(10.0), Inches(0.7), "AIPIs helps teams ship multi-model products faster while retaining visibility and flexibility as the market changes.", size=15, color=MUTED)


def market(slide):
    add_bg(slide, dark=False)
    add_tag(slide, Inches(0.72), Inches(0.5), "Market", dark=False)
    add_text(slide, Inches(0.72), Inches(1.0), Inches(7.8), Inches(0.95), "Customers range from individual power users to developer platforms.", size=24, color=BLACK, bold=True)
    add_card(slide, Inches(0.72), Inches(3.1), Inches(2.75), Inches(1.95), "Power users", "Users who want access to multiple AI models without managing separate tools.", CORAL)
    add_card(slide, Inches(3.62), Inches(3.1), Inches(2.75), Inches(1.95), "Startups", "Product teams building features on top of more than one provider.", CYAN)
    add_card(slide, Inches(6.52), Inches(3.1), Inches(2.75), Inches(1.95), "Enterprise teams", "Operators who need governance, routing, cost control, and provider flexibility.", INDIGO)
    add_card(slide, Inches(9.42), Inches(3.1), Inches(2.9), Inches(1.95), "Developers", "Teams who want one API instead of many provider contracts and integrations.", LIME)


def business(slide):
    add_bg(slide, dark=True)
    add_tag(slide, Inches(0.72), Inches(0.5), "Business model", dark=True)
    add_text(slide, Inches(0.72), Inches(1.0), Inches(8.0), Inches(0.95), "A product business with platform and API upside.", size=24, color=WHITE, bold=True)
    add_card(slide, Inches(0.72), Inches(3.2), Inches(3.45), Inches(1.95), "Subscription access", "Users pay for premium model access, higher limits, and provider-specific features.", CORAL, dark=True)
    add_card(slide, Inches(4.42), Inches(3.2), Inches(3.45), Inches(1.95), "B2B platform", "Companies license orchestration and control capabilities for internal and product use.", CYAN, dark=True)
    add_card(slide, Inches(8.12), Inches(3.2), Inches(4.15), Inches(1.95), "Unified AI API", "External developers pay for AIPIs endpoints rather than integrating every provider themselves.", INDIGO, dark=True)


def roadmap(slide):
    add_bg(slide, dark=False)
    add_tag(slide, Inches(0.72), Inches(0.5), "Roadmap", dark=False)
    add_text(slide, Inches(0.72), Inches(1.0), Inches(8.5), Inches(0.95), "Start with provider access, expand into orchestration infrastructure.", size=24, color=BLACK, bold=True)
    add_card(slide, Inches(0.72), Inches(3.2), Inches(3.5), Inches(1.95), "Phase 1", "Chat UI, initial providers, unified request layer, and account-level usage controls.", CYAN)
    add_card(slide, Inches(4.45), Inches(3.2), Inches(3.5), Inches(1.95), "Phase 2", "Routing logic, fallback rules, pricing tiers, and observability dashboards.", INDIGO)
    add_card(slide, Inches(8.18), Inches(3.2), Inches(3.5), Inches(1.95), "Phase 3", "Enterprise policy controls, workflow APIs, and deeper private model support.", CORAL)
    add_text(slide, Inches(0.72), Inches(5.75), Inches(10.2), Inches(0.7), "The long-term position is a default access layer for teams that want provider agility without operational fragmentation.", size=15, color=SOFT)


def closing(slide):
    add_bg(slide, dark=True)
    add_orb(slide, Inches(9.8), Inches(4.7), Inches(2.2), CYAN, 0)
    add_orb(slide, Inches(10.55), Inches(4.9), Inches(1.55), INDIGO, 0)
    add_orb(slide, Inches(11.18), Inches(5.08), Inches(1.05), CORAL, 0)
    add_tag(slide, Inches(0.72), Inches(0.5), "Closing", dark=True)
    add_text(slide, Inches(0.72), Inches(1.12), Inches(7.9), Inches(1.25), "AIPIs is building the access layer for the multi-model AI era.", size=29, color=WHITE, bold=True)
    add_text(slide, Inches(0.72), Inches(2.6), Inches(6.8), Inches(1.3), "One interface, one gateway, and one control plane for the AI providers people actually want to use.", size=19, color=MUTED)
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(5.55), Inches(4.8), Inches(0.95))
    box.fill.solid()
    box.fill.fore_color.rgb = DEEP
    box.line.fill.background()
    add_text(slide, Inches(0.95), Inches(5.75), Inches(1.0), Inches(0.2), "Contact", size=11, color=LIME, bold=True)
    add_text(slide, Inches(0.95), Inches(6.0), Inches(4.0), Inches(0.28), "Founding platform deck for AIPIs", size=16, color=WHITE, bold=True)
    add_text(slide, Inches(0.95), Inches(6.25), Inches(3.5), Inches(0.2), "multi-provider ai access infrastructure", size=11, color=MUTED)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]

    for fn in [cover, problem, product, architecture, moat, market, business, roadmap, closing]:
        slide = prs.slides.add_slide(layout)
        fn(slide)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
